import os
import time

class DocumentFormats:
    @staticmethod
    def convertir(origen_path: str, destino_path: str, from_ext: str, to_ext: str) -> str:
        """
        Enruta la conversión a la función correspondiente según el formato origen y destino.
        """
        origen_path = os.path.abspath(origen_path)
        destino_path = os.path.abspath(destino_path)
        
        try:
            if from_ext == "jpg" and to_ext == "pdf":
                return DocumentFormats.jpg_to_pdf(origen_path, destino_path)
            elif from_ext == "pdf" and to_ext == "jpg":
                return DocumentFormats.pdf_to_jpg(origen_path, destino_path)
            elif from_ext == "pdf" and to_ext == "docx":
                return DocumentFormats.pdf_to_word(origen_path, destino_path)
            elif from_ext == "docx" and to_ext == "pdf":
                return DocumentFormats.word_to_pdf(origen_path, destino_path)
            elif from_ext == "xlsx" and to_ext == "pdf":
                return DocumentFormats.excel_to_pdf(origen_path, destino_path)
            elif from_ext == "pptx" and to_ext == "pdf":
                return DocumentFormats.powerpoint_to_pdf(origen_path, destino_path)
            elif from_ext == "pdf" and to_ext == "xlsx":
                return DocumentFormats.pdf_to_excel(origen_path, destino_path)
            elif from_ext == "pdf" and to_ext == "pptx":
                return DocumentFormats.pdf_to_powerpoint(origen_path, destino_path)
            elif from_ext == "html" and to_ext == "pdf":
                return DocumentFormats.html_to_pdf(origen_path, destino_path)
            else:
                return f"❌ Conversión de {from_ext.upper()} a {to_ext.upper()} no está soportada aún."
        except Exception as e:
            return f"❌ Error al convertir: {e}"

    @staticmethod
    def jpg_to_pdf(origen, destino):
        import img2pdf
        with open(destino, "wb") as f:
            f.write(img2pdf.convert(origen))
        return "✅ Archivo JPG convertido a PDF exitosamente."

    @staticmethod
    def pdf_to_jpg(origen, destino):
        import fitz
        import os
        
        doc = fitz.open(origen)
        num_pages = len(doc)
        
        if num_pages == 0:
            return "❌ El PDF está vacío."
            
        output_dir = os.path.dirname(destino)
        base_name = os.path.splitext(os.path.basename(destino))[0]
        
        # Si tiene más de 10 hojas, creamos una carpeta dedicada
        if num_pages > 5:
            target_dir = os.path.join(output_dir, base_name)
            if not os.path.exists(target_dir):
                os.makedirs(target_dir)
        else:
            target_dir = output_dir
            
        # Extraer todas las páginas
        for i in range(num_pages):
            page = doc.load_page(i)
            pix = page.get_pixmap(dpi=300)
            
            # Si es solo 1 hoja, mantiene el nombre original, sino enumera
            if num_pages == 1:
                file_path = os.path.join(target_dir, f"{base_name}.jpg")
            else:
                file_path = os.path.join(target_dir, f"{base_name}_pagina_{i+1}.jpg")
                
            pix.save(file_path)
            
        doc.close()
        
        if num_pages > 5:
            return f"✅ {num_pages} páginas extraídas a JPG y guardadas en la carpeta '{base_name}'."
        elif num_pages > 1:
            return f"✅ {num_pages} páginas extraídas a JPG exitosamente."
        else:
            return "✅ Página convertida a JPG exitosamente."

    @staticmethod
    def pdf_to_word(origen, destino):
        from pdf2docx import Converter
        cv = Converter(origen)
        cv.convert(destino)
        cv.close()
        return "✅ Archivo PDF convertido a Word (DOCX) exitosamente."

    @staticmethod
    def word_to_pdf(origen, destino):
        from docx2pdf import convert
        convert(origen, destino)
        return "✅ Archivo Word convertido a PDF exitosamente."

    @staticmethod
    def excel_to_pdf(origen, destino):
        import comtypes.client
        try:
            excel = comtypes.client.CreateObject("Excel.Application")
            excel.Visible = False
            wb = excel.Workbooks.Open(origen)
            # 0 = xlTypePDF
            wb.ExportAsFixedFormat(0, destino)
            wb.Close()
            excel.Quit()
            return "✅ Archivo Excel convertido a PDF exitosamente."
        except Exception as e:
            return f"❌ Error usando Excel (comtypes): {e}. Asegúrate de tener MS Excel instalado."

    @staticmethod
    def powerpoint_to_pdf(origen, destino):
        import comtypes.client
        try:
            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
            # 32 = ppFixedFormatTypePDF
            presentation = powerpoint.Presentations.Open(origen, WithWindow=False)
            presentation.ExportAsFixedFormat(destino, 32)
            presentation.Close()
            powerpoint.Quit()
            return "✅ Archivo PowerPoint convertido a PDF exitosamente."
        except Exception as e:
            return f"❌ Error usando PowerPoint (comtypes): {e}. Asegúrate de tener MS PowerPoint instalado."

    @staticmethod
    def pdf_to_excel(origen, destino):
        import pdfplumber
        import pandas as pd
        with pdfplumber.open(origen) as pdf:
            all_tables = []
            for page in pdf.pages:
                tables = page.extract_tables()
                for table in tables:
                    if not table:
                        continue
                    df = pd.DataFrame(table[1:], columns=table[0])
                    all_tables.append(df)
            
            if all_tables:
                with pd.ExcelWriter(destino) as writer:
                    for i, df in enumerate(all_tables):
                        df.to_excel(writer, sheet_name=f"Tabla_{i+1}", index=False)
                return "✅ Tablas del PDF extraídas a Excel exitosamente."
            else:
                # Crear un excel vacío si no hay tablas
                pd.DataFrame(["No se encontraron tablas"]).to_excel(destino, index=False)
                return "✅ Archivo Excel creado, pero no se detectaron tablas claras en el PDF."

    @staticmethod
    def pdf_to_powerpoint(origen, destino):
        import fitz
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6] # Blank
        
        doc = fitz.open(origen)
        for i in range(len(doc)):
            page = doc.load_page(i)
            pix = page.get_pixmap(dpi=150)
            img_path = f"temp_page_{i}_{int(time.time())}.png"
            pix.save(img_path)
            
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
            
            if os.path.exists(img_path):
                os.remove(img_path)
        
        prs.save(destino)
        doc.close()
        return "✅ Páginas del PDF convertidas a diapositivas de PowerPoint exitosamente."

    @staticmethod
    def html_to_pdf(origen, destino):
        from PyQt6.QtGui import QTextDocument
        from PyQt6.QtPrintSupport import QPrinter
        try:
            with open(origen, "r", encoding="utf-8") as f:
                html_content = f.read()
                
            doc = QTextDocument()
            doc.setHtml(html_content)
            
            printer = QPrinter(QPrinter.PrinterMode.HighResolution)
            printer.setOutputFormat(QPrinter.OutputFormat.PdfFormat)
            printer.setOutputFileName(destino)
            
            doc.print(printer)
            return "✅ Archivo HTML convertido a PDF exitosamente."
        except Exception as e:
            return f"❌ Error al convertir HTML: {e}"

